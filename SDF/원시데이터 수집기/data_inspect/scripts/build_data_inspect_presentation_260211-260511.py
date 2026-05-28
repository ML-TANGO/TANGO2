#!/usr/bin/env python3
"""Build a presentation deck from data_inspect_260211-260511 notebook outputs."""

from __future__ import annotations

from pathlib import Path
from typing import Iterable, Optional

import matplotlib

matplotlib.use("Agg")

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[2]
OUTPUT_ROOT = PROJECT_ROOT / "data_inspect" / "output"
BASE_OUTPUT = OUTPUT_ROOT / "data_inspect_260211-260511"
B_OUTPUT = OUTPUT_ROOT / "motor_temperature_daily_pattern_260211-260511"
C_OUTPUT = OUTPUT_ROOT / "weekly_motor_weather_analysis_260211-260511"
D_OUTPUT = OUTPUT_ROOT / "weekly_motor_humidity_analysis_260211-260511"
PRESENTATION_ROOT = OUTPUT_ROOT / "presentation_260211-260511"
PLOT_ROOT = PRESENTATION_ROOT / "plots"
PPTX_PATH = PRESENTATION_ROOT / "data_inspect_260211-260511_presentation.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT = "AppleGothic"
DARK = RGBColor(32, 43, 54)
BLUE = RGBColor(31, 119, 180)
GREEN = RGBColor(44, 160, 44)
ORANGE = RGBColor(255, 127, 14)
RED = RGBColor(214, 39, 40)
GRAY = RGBColor(105, 117, 130)
LIGHT_BG = RGBColor(245, 247, 250)


PRESENTATION_ROOT.mkdir(parents=True, exist_ok=True)
PLOT_ROOT.mkdir(parents=True, exist_ok=True)

plt.rcParams["axes.unicode_minus"] = False
for font_name in ["AppleGothic", "Malgun Gothic", "NanumGothic", "DejaVu Sans"]:
    plt.rcParams["font.family"] = font_name
    break


def read_csv(path: Path) -> pd.DataFrame:
    return pd.read_csv(path) if path.exists() else pd.DataFrame()


def add_textbox(slide, text: str, left, top, width, height, font_size=18, bold=False, color=DARK, align=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    if align is not None:
        paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title: str, subtitle: str = ""):
    add_textbox(slide, title, Inches(0.55), Inches(0.25), Inches(12.2), Inches(0.55), 27, True, DARK)
    if subtitle:
        add_textbox(slide, subtitle, Inches(0.58), Inches(0.83), Inches(12.1), Inches(0.35), 12, False, GRAY)
    line = slide.shapes.add_shape(1, Inches(0.55), Inches(1.17), Inches(12.2), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(220, 226, 232)
    line.line.color.rgb = RGBColor(220, 226, 232)


def add_bullets(slide, items: Iterable[str], left, top, width, height, font_size=17, color=DARK):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.font.name = FONT
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = color
        paragraph.space_after = Pt(8)
    return box


def add_metric_card(slide, title: str, value: str, note: str, left, top, width, height, accent=BLUE):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.color.rgb = RGBColor(222, 228, 235)
    accent_shape = slide.shapes.add_shape(1, left, top, Inches(0.08), height)
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = accent
    accent_shape.line.color.rgb = accent
    add_textbox(slide, title, left + Inches(0.18), top + Inches(0.12), width - Inches(0.25), Inches(0.28), 11, True, GRAY)
    add_textbox(slide, value, left + Inches(0.18), top + Inches(0.45), width - Inches(0.25), Inches(0.45), 24, True, DARK)
    add_textbox(slide, note, left + Inches(0.18), top + Inches(0.98), width - Inches(0.25), height - Inches(1.05), 10, False, GRAY)


def add_image(slide, path: Optional[Path], left, top, width, height):
    if path is not None and path.exists():
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    else:
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(250, 250, 250)
        shape.line.color.rgb = RGBColor(220, 220, 220)
        add_textbox(slide, "이미지 없음", left, top + height / 2 - Inches(0.15), width, Inches(0.3), 14, False, GRAY, PP_ALIGN.CENTER)


def add_table(slide, dataframe: pd.DataFrame, left, top, width, height, font_size=8):
    if dataframe.empty:
        add_textbox(slide, "표 데이터 없음", left, top, width, height, 12, False, GRAY)
        return
    rows = min(len(dataframe), 8) + 1
    cols = min(len(dataframe.columns), 5)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    for col_idx, column in enumerate(dataframe.columns[:cols]):
        cell = table.cell(0, col_idx)
        cell.text = str(column)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(231, 237, 244)
    for row_idx, row in enumerate(dataframe.iloc[: rows - 1, :cols].itertuples(index=False), start=1):
        for col_idx, value in enumerate(row):
            text = "" if pd.isna(value) else str(value)
            if isinstance(value, float):
                text = f"{value:.3f}" if abs(value) < 10 else f"{value:.1f}"
            table.cell(row_idx, col_idx).text = text[:38]
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = FONT
                paragraph.font.size = Pt(font_size)


def build_hourly_join() -> tuple[pd.DataFrame, pd.DataFrame, pd.DataFrame]:
    open_rate = read_csv(B_OUTPUT / "open_rate_hourly_long.csv")
    temp = read_csv(B_OUTPUT / "temperature_hourly_long.csv")
    if open_rate.empty or temp.empty:
        return pd.DataFrame(), pd.DataFrame(), pd.DataFrame()
    open_rate["local_ts"] = pd.to_datetime(open_rate["local_ts"], errors="coerce")
    temp["local_ts"] = pd.to_datetime(temp["local_ts"], errors="coerce")
    motor_hourly = open_rate.groupby("local_ts", as_index=False).agg(
        mean_open_rate=("open_rate", "mean"),
        median_open_rate=("open_rate", "median"),
        high_open_channel_ratio=("open_rate", lambda x: float((x >= 90).mean())),
        channel_count=("open_rate", "count"),
    )
    internal = temp[temp["scope"] == "internal"].groupby("local_ts", as_index=False).agg(internal_temp=("temp", "mean"))
    external = temp[temp["scope"] == "external"].groupby("local_ts", as_index=False).agg(external_temp=("temp", "mean"))
    joined = motor_hourly.merge(internal, on="local_ts", how="inner").merge(external, on="local_ts", how="left")
    joined["month"] = joined["local_ts"].dt.strftime("%Y-%m")
    return joined, open_rate, temp


def calculate_summary() -> dict:
    file_overview = read_csv(BASE_OUTPUT / "file_overview.csv")
    gaps = read_csv(BASE_OUTPUT / "gaps_over_5h.csv")
    delays = read_csv(BASE_OUTPUT / "ts_created_delay_over_3h.csv")
    base_corr = read_csv(BASE_OUTPUT / "open_rate_environment_correlations.csv")
    daily_open = read_csv(B_OUTPUT / "daily_open_rate_summary.csv")
    daily_temp = read_csv(B_OUTPUT / "daily_temperature_summary.csv")
    weather_corr = read_csv(C_OUTPUT / "all_week_correlations.csv")
    weather_daily = read_csv(C_OUTPUT / "all_week_daily_summary.csv")
    humidity_corr = read_csv(D_OUTPUT / "all_week_humidity_correlations.csv")
    humidity_lag = read_csv(D_OUTPUT / "all_week_humidity_lag_correlations.csv")
    humidity_daily = read_csv(D_OUTPUT / "all_week_daily_humidity_summary.csv")
    joined, open_rate, temp = build_hourly_join()

    internal = temp[temp["scope"] == "internal"] if not temp.empty else pd.DataFrame()
    feb_mar = joined[joined["month"].isin(["2026-02", "2026-03"])] if not joined.empty else pd.DataFrame()
    candidate = feb_mar[(feb_mar["mean_open_rate"] >= 5) & (feb_mar["mean_open_rate"] <= 60)] if not feb_mar.empty else pd.DataFrame()
    if candidate.empty:
        candidate = feb_mar

    monthly = pd.DataFrame()
    if not joined.empty:
        monthly = joined.groupby("month").agg(
            mean_open_rate=("mean_open_rate", "mean"),
            median_open_rate=("median_open_rate", "median"),
            high_open_channel_ratio=("high_open_channel_ratio", "mean"),
            internal_temp_mean=("internal_temp", "mean"),
            internal_temp_min=("internal_temp", "min"),
            internal_temp_max=("internal_temp", "max"),
        ).reset_index()
        monthly["open_temp_corr"] = joined.groupby("month").apply(lambda x: x["mean_open_rate"].corr(x["internal_temp"])).values

    temp_range = {
        "observed_min": float(internal["temp"].min()) if not internal.empty else np.nan,
        "observed_max": float(internal["temp"].max()) if not internal.empty else np.nan,
        "observed_p10": float(internal["temp"].quantile(0.10)) if not internal.empty else np.nan,
        "observed_p50": float(internal["temp"].quantile(0.50)) if not internal.empty else np.nan,
        "observed_p90": float(internal["temp"].quantile(0.90)) if not internal.empty else np.nan,
        "candidate_min": float(candidate["internal_temp"].quantile(0.10)) if not candidate.empty else np.nan,
        "candidate_max": float(candidate["internal_temp"].quantile(0.90)) if not candidate.empty else np.nan,
        "candidate_optimal": float(candidate["internal_temp"].median()) if not candidate.empty else np.nan,
        "candidate_rows": int(len(candidate)),
    }

    april_may = open_rate[open_rate["month"].isin(["2026-04", "2026-05"])] if not open_rate.empty else pd.DataFrame()
    april_may_high_ratio = float((april_may["open_rate"] >= 90).mean()) if not april_may.empty else np.nan
    april_may_mean_open = float(april_may["open_rate"].mean()) if not april_may.empty else np.nan

    return {
        "file_overview": file_overview,
        "gaps": gaps,
        "delays": delays,
        "base_corr": base_corr,
        "daily_open": daily_open,
        "daily_temp": daily_temp,
        "weather_corr": weather_corr,
        "weather_daily": weather_daily,
        "humidity_corr": humidity_corr,
        "humidity_lag": humidity_lag,
        "humidity_daily": humidity_daily,
        "joined": joined,
        "open_rate": open_rate,
        "temp": temp,
        "monthly": monthly,
        "temp_range": temp_range,
        "april_may_high_ratio": april_may_high_ratio,
        "april_may_mean_open": april_may_mean_open,
    }


def create_summary_plots(summary: dict) -> dict[str, Path]:
    plots: dict[str, Path] = {}
    joined = summary["joined"]
    monthly = summary["monthly"]
    temp = summary["temp"]

    if not monthly.empty:
        path = PLOT_ROOT / "monthly_open_rate_temperature_summary.png"
        fig, ax1 = plt.subplots(figsize=(11, 5.8))
        x = np.arange(len(monthly))
        ax1.bar(x - 0.18, monthly["mean_open_rate"], width=0.36, color="#1f77b4", alpha=0.8, label="평균 개폐율")
        ax1.bar(x + 0.18, monthly["high_open_channel_ratio"] * 100, width=0.36, color="#ff7f0e", alpha=0.8, label="90% 이상 개방 채널 비율")
        ax1.set_ylabel("개폐율 / 비율(%)")
        ax1.set_ylim(0, 105)
        ax1.set_xticks(x)
        ax1.set_xticklabels(monthly["month"])
        ax2 = ax1.twinx()
        ax2.plot(x, monthly["internal_temp_mean"], color="#2ca02c", marker="o", linewidth=2.2, label="평균 내부온도")
        ax2.set_ylabel("내부온도(°C)")
        lines1, labels1 = ax1.get_legend_handles_labels()
        lines2, labels2 = ax2.get_legend_handles_labels()
        ax1.legend(lines1 + lines2, labels1 + labels2, loc="upper left")
        ax1.set_title("월별 개폐율과 내부온도 요약")
        fig.tight_layout()
        fig.savefig(path, dpi=160)
        plt.close(fig)
        plots["monthly_summary"] = path

    if not joined.empty:
        path = PLOT_ROOT / "open_rate_vs_internal_temp_scatter.png"
        fig, ax = plt.subplots(figsize=(10.5, 5.8))
        for month, group in joined.groupby("month"):
            ax.scatter(group["internal_temp"], group["mean_open_rate"], s=18, alpha=0.6, label=month)
        ax.set_xlabel("내부온도(°C)")
        ax.set_ylabel("채널 평균 개폐율(%)")
        ax.set_ylim(-5, 105)
        ax.set_title("내부온도와 개폐율 산점도: 온도가 높을수록 개폐율도 높은 경향")
        ax.legend(loc="upper left", ncol=2, fontsize=8)
        fig.tight_layout()
        fig.savefig(path, dpi=160)
        plt.close(fig)
        plots["scatter"] = path

    if not temp.empty:
        internal = temp[temp["scope"] == "internal"].copy()
        if not internal.empty:
            internal["month"] = pd.to_datetime(internal["local_ts"], errors="coerce").dt.strftime("%Y-%m")
            path = PLOT_ROOT / "internal_temperature_boxplot_by_month.png"
            months = sorted(internal["month"].dropna().unique())
            data = [internal.loc[internal["month"] == month, "temp"].dropna() for month in months]
            fig, ax = plt.subplots(figsize=(9.8, 5.6))
            ax.boxplot(data, labels=months, showfliers=False)
            ax.set_ylabel("내부온도(°C)")
            ax.set_title("월별 내부온도 분포")
            fig.tight_layout()
            fig.savefig(path, dpi=160)
            plt.close(fig)
            plots["temp_boxplot"] = path

    summary_table = PLOT_ROOT / "temperature_candidate_table.csv"
    temp_range = summary["temp_range"]
    pd.DataFrame(
        [
            {"항목": "전체 관측 내부온도 최소", "값": f"{temp_range['observed_min']:.2f}°C"},
            {"항목": "전체 관측 내부온도 최대", "값": f"{temp_range['observed_max']:.2f}°C"},
            {"항목": "전체 관측 p10~p90", "값": f"{temp_range['observed_p10']:.2f}~{temp_range['observed_p90']:.2f}°C"},
            {"항목": "2~3월 제어 후보 범위", "값": f"{temp_range['candidate_min']:.2f}~{temp_range['candidate_max']:.2f}°C"},
            {"항목": "최적온도 후보 중앙값", "값": f"{temp_range['candidate_optimal']:.2f}°C"},
        ]
    ).to_csv(summary_table, index=False, encoding="utf-8-sig")
    plots["temperature_candidate_table"] = summary_table
    return plots


def make_deck(summary: dict, plots: dict[str, Path]) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(243, 247, 251)
    shape.line.color.rgb = RGBColor(243, 247, 251)
    add_textbox(slide, "환경정보 기반 목표 개폐율 학습을 위한\n데이터 분석 결과", Inches(0.75), Inches(1.35), Inches(12), Inches(1.25), 32, True, DARK)
    add_textbox(slide, "2026-02-11 ~ 2026-05-11 센서·기상청 데이터 분석", Inches(0.78), Inches(2.75), Inches(10), Inches(0.42), 17, False, GRAY)
    add_metric_card(slide, "분석 노트북", "A/B/C/D", "데이터 품질, 일별 패턴, 주 단위 기상·습도 분석", Inches(0.8), Inches(4.35), Inches(3.75), Inches(1.45), BLUE)
    add_metric_card(slide, "모터 채널", "16개", "단동 7동/8동, 각 8채널 개별 분석", Inches(4.85), Inches(4.35), Inches(3.75), Inches(1.45), GREEN)
    add_metric_card(slide, "핵심 이슈", "4~5월 상시 개방", "고개폐율이 내부온도 상승과 함께 관측되어 상세 분석 필요", Inches(8.9), Inches(4.35), Inches(3.75), Inches(1.45), ORANGE)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "1. 발표 목적과 분석 방향", "원래 목표는 ML 모델 학습용 온도 기준 추출이었으나, 데이터 특성상 상세 검증이 필요했다.")
    add_bullets(
        slide,
        [
            "원래 목적: 환경정보를 입력으로 목표 개폐율을 출력하는 모델을 학습시키기 위한 내부온도 최소·최대·최적값 후보 추출",
            "기대 가정: 내부온도가 높아지면 개폐율이 증가하고, 적정 온도 구간에서는 개폐율이 안정화될 것",
            "확인된 문제: 4~5월에는 개폐율이 거의 항상 높고, 높은 개폐율일수록 내부온도도 높게 나타나는 양의 상관 가능성",
            "따라서 결론을 바로 온도 기준으로 확정하지 않고, 시간대·월별·주 단위·습도/강우/일사까지 상세 분석했다.",
        ],
        Inches(0.75), Inches(1.55), Inches(12), Inches(4.5), 18,
    )

    file_count = int(len(summary["file_overview"])) if not summary["file_overview"].empty else 0
    gap_count = int(len(summary["gaps"])) if not summary["gaps"].empty else 0
    delay_count = int(len(summary["delays"])) if not summary["delays"].empty else 0
    corr_count = int(len(summary["base_corr"])) if not summary["base_corr"].empty else 0

    slide = prs.slides.add_slide(blank)
    add_title(slide, "2. 사용 데이터와 산출물 구성", "각 notebook 결과를 하나의 발표 흐름으로 연결")
    add_metric_card(slide, "분석 파일 수", f"{file_count}개", "organized_by_sensor CSV 기준", Inches(0.7), Inches(1.45), Inches(2.8), Inches(1.25), BLUE)
    add_metric_card(slide, "5시간 이상 결측", f"{gap_count}구간", "장기 결측은 단순 보간 위험", Inches(3.8), Inches(1.45), Inches(2.8), Inches(1.25), RED)
    add_metric_card(slide, "3시간 이상 지연", f"{delay_count:,}행", "ts와 created_at 차이", Inches(6.9), Inches(1.45), Inches(2.8), Inches(1.25), ORANGE)
    add_metric_card(slide, "상관 계산", f"{corr_count}건", "개폐율-환경 변수 5분 리샘플링", Inches(10.0), Inches(1.45), Inches(2.8), Inches(1.25), GREEN)
    add_bullets(
        slide,
        [
            "기본 노트북: 데이터 품질, 결측, 지연, 개폐율-환경 상관 분석",
            "B 노트북: 모터 개폐 이벤트, 시간대별 패턴, 실제 open_rate와 내부/외부온도 변화",
            "C 노트북/스크립트: 주 단위 open_rate·온도·강우·일사/일조 분석",
            "D 노트북/스크립트: 주 단위 open_rate·내외부/KMA 습도 및 lag 상관 분석",
        ],
        Inches(0.85), Inches(3.05), Inches(11.7), Inches(3.5), 16,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "3. 내부온도 기준 추출 결과와 한계", "관측 범위는 산출 가능하지만, 최적온도는 현재 데이터만으로 확정하기 어렵다.")
    t = summary["temp_range"]
    add_metric_card(slide, "관측 내부온도 최소", f"{t['observed_min']:.1f}°C", "전체 내부 온습도 센서 기준", Inches(0.65), Inches(1.45), Inches(3.0), Inches(1.25), BLUE)
    add_metric_card(slide, "관측 내부온도 최대", f"{t['observed_max']:.1f}°C", "전체 내부 온습도 센서 기준", Inches(3.95), Inches(1.45), Inches(3.0), Inches(1.25), RED)
    add_metric_card(slide, "2~3월 후보 범위", f"{t['candidate_min']:.1f}~{t['candidate_max']:.1f}°C", "중간 개폐율 구간 중심 p10~p90", Inches(7.25), Inches(1.45), Inches(3.0), Inches(1.25), GREEN)
    add_metric_card(slide, "최적온도 후보", f"{t['candidate_optimal']:.1f}°C", "확정값이 아닌 관측 기반 중앙값", Inches(10.55), Inches(1.45), Inches(2.15), Inches(1.25), ORANGE)
    add_image(slide, plots.get("temp_boxplot"), Inches(0.75), Inches(3.05), Inches(5.8), Inches(3.55))
    add_bullets(
        slide,
        [
            "4~5월에는 이미 많은 채널이 높은 개폐율 상태라서 ‘개폐율이 높으면 내부온도가 내려간다’는 형태의 학습 신호가 약하다.",
            "오히려 고개폐율과 고내부온도가 함께 관측되어, 모델이 개폐율 증가를 온도 상승의 결과처럼 학습할 위험이 있다.",
            "따라서 최적온도는 단일 숫자로 확정하지 않고, 2~3월 중심 후보 범위와 운영 지식 기반 검증이 필요하다.",
        ],
        Inches(6.8), Inches(3.1), Inches(5.8), Inches(3.2), 15,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "4. 4~5월 데이터의 핵심 문제", "개폐율이 거의 항상 높고, 내부온도도 함께 높아지는 구조")
    add_metric_card(slide, "4~5월 평균 개폐율", f"{summary['april_may_mean_open']:.1f}%", "채널-hour 샘플 평균", Inches(0.8), Inches(1.45), Inches(3.2), Inches(1.25), BLUE)
    add_metric_card(slide, "4~5월 90% 이상 개방 비율", f"{summary['april_may_high_ratio'] * 100:.1f}%", "채널-hour 샘플 기준", Inches(4.25), Inches(1.45), Inches(3.4), Inches(1.25), ORANGE)
    if not summary["monthly"].empty:
        corr_text = ", ".join(f"{row.month}: {row.open_temp_corr:.2f}" for row in summary["monthly"].itertuples(index=False))
    else:
        corr_text = "계산 불가"
    add_metric_card(slide, "월별 온도-개폐율 상관", "양의 상관 주의", corr_text[:95], Inches(7.9), Inches(1.45), Inches(4.55), Inches(1.25), RED)
    add_image(slide, plots.get("monthly_summary"), Inches(0.75), Inches(3.0), Inches(5.95), Inches(3.65))
    add_image(slide, plots.get("scatter"), Inches(6.95), Inches(3.0), Inches(5.65), Inches(3.65))

    slide = prs.slides.add_slide(blank)
    add_title(slide, "5. B 노트북: 실제 open_rate와 온도 추세", "이벤트 수만이 아니라 실제 개폐율 값의 전체 기간 변화를 확인")
    add_image(slide, B_OUTPUT / "plots" / "daily_open_rate_temperature_trend.png", Inches(0.7), Inches(1.45), Inches(5.95), Inches(4.85))
    add_image(slide, B_OUTPUT / "plots" / "daily_open_rate_by_channel_all.png", Inches(6.85), Inches(1.45), Inches(5.8), Inches(4.85))
    add_textbox(slide, "핵심 메시지: 4~5월로 갈수록 개폐율이 높은 상태가 장시간 유지되어, 단순 상관 기반 최적온도 추출은 왜곡될 수 있다.", Inches(0.9), Inches(6.45), Inches(11.7), Inches(0.4), 15, True, RED, PP_ALIGN.CENTER)

    top_base = summary["base_corr"].copy()
    if not top_base.empty:
        top_base = top_base.sort_values("abs_pearson_corr", ascending=False).head(6)
        top_base = top_base[["motor_metric", "environment_metric", "pearson_corr", "spearman_corr", "best_lag_minutes_by_abs_corr"]]
        top_base.columns = ["개폐율", "환경변수", "Pearson", "Spearman", "lag(분)"]

    slide = prs.slides.add_slide(blank)
    add_title(slide, "6. 기본 상관 분석: 온도가 가장 강한 신호로 관측", "5분 리샘플링 기준 개폐율-환경변수 상관")
    add_image(slide, BASE_OUTPUT / "plots" / "open_rate_environment_correlation_heatmap.png", Inches(0.7), Inches(1.45), Inches(5.85), Inches(4.8))
    add_table(slide, top_base, Inches(6.8), Inches(1.45), Inches(5.85), Inches(3.0), 7)
    add_bullets(
        slide,
        [
            "상위 상관은 주로 외부 온습도 센서의 온도 변수에서 나타났다.",
            "다만 상관은 인과가 아니며, 4~5월 상시 개방 구조 때문에 목표 개폐율 라벨로 바로 쓰기 어렵다.",
            "모델 학습 시 lag, rolling, 시간대, 일사/강우, 채널 식별 feature를 함께 사용해야 한다.",
        ],
        Inches(6.85), Inches(4.65), Inches(5.75), Inches(1.8), 13,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "7. C 노트북/스크립트: 주 단위 기상·온도 상세 분석", "전체 기간을 한 번에 보지 않고 7일 단위로 분해")
    add_image(slide, C_OUTPUT / "week_10_20260422" / "week_overview_open_rate_temperature_rain_solar.png", Inches(0.65), Inches(1.35), Inches(6.1), Inches(5.45))
    add_image(slide, C_OUTPUT / "week_10_20260422" / "week_correlation_heatmap.png", Inches(6.95), Inches(1.35), Inches(5.65), Inches(5.45))
    add_textbox(slide, "대표 예시: 2026-04-22 주차. 개별 채널 open_rate와 내부/외부온도, 강우, 일사/일조를 같은 축에서 확인.", Inches(0.85), Inches(6.85), Inches(11.6), Inches(0.35), 12, False, GRAY, PP_ALIGN.CENTER)

    top_weather = summary["weather_corr"].copy()
    if not top_weather.empty:
        top_weather["abs_pearson"] = top_weather["pearson"].abs()
        top_weather = top_weather.sort_values("abs_pearson", ascending=False).head(7)
        top_weather = top_weather[["week_index", "motor_channel", "weather_variable", "pearson", "spearman"]]
        top_weather.columns = ["주", "채널", "변수", "Pearson", "Spearman"]

    slide = prs.slides.add_slide(blank)
    add_title(slide, "8. 주 단위 기상 분석 요약", "채널별로 온도·강우·일사 변수와의 관련성이 다르게 나타남")
    add_table(slide, top_weather, Inches(0.75), Inches(1.45), Inches(6.2), Inches(4.1), 7)
    add_bullets(
        slide,
        [
            "평균 개폐율 하나로 합치지 않고 16개 채널을 각각 분석했다.",
            "채널별 위치와 제어 특성 차이 때문에 같은 환경 변수에도 상관 크기가 다르다.",
            "일사/일조는 온도 상승과 함께 작용하므로, 단일 온도 기준만으로 목표 개폐율을 정의하면 누락 변수가 생긴다.",
            "주 단위 산출물은 스크립트로 전체 13주를 일괄 생성하도록 구성했다.",
        ],
        Inches(7.25), Inches(1.55), Inches(5.3), Inches(4.7), 15,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "9. D 노트북/스크립트: 습도와 개폐율 관계", "내부/외부/KMA 습도와 lag 후보를 별도로 확인")
    add_image(slide, D_OUTPUT / "week_10_20260422" / "week_open_rate_humidity_overview.png", Inches(0.65), Inches(1.35), Inches(6.1), Inches(5.45))
    add_image(slide, D_OUTPUT / "week_10_20260422" / "week_humidity_lag_top.png", Inches(6.95), Inches(1.35), Inches(5.65), Inches(5.45))
    add_textbox(slide, "습도는 강우·온도와 동시 변동하므로 독립 원인으로 확정하지 않고, lag/rolling feature 후보로 활용한다.", Inches(0.85), Inches(6.85), Inches(11.6), Inches(0.35), 12, False, GRAY, PP_ALIGN.CENTER)

    top_humidity = summary["humidity_corr"].copy()
    if not top_humidity.empty:
        top_humidity = top_humidity.sort_values("abs_pearson", ascending=False).head(7)
        top_humidity = top_humidity[["week_index", "motor_channel", "humidity_variable", "pearson", "spearman"]]
        top_humidity.columns = ["주", "채널", "습도변수", "Pearson", "Spearman"]

    slide = prs.slides.add_slide(blank)
    add_title(slide, "10. 습도 분석 요약", "습도 자체보다 온도·강우와 결합한 feature로 보는 것이 안전")
    add_table(slide, top_humidity, Inches(0.75), Inches(1.45), Inches(6.3), Inches(4.25), 7)
    add_bullets(
        slide,
        [
            "동시 상관과 lag 상관을 모두 계산해 습도 변화 선행/후행 후보를 확인했다.",
            "고습도 구간은 강우·외기 변화와 결합되어 나타나므로 인과 해석은 제한적이다.",
            "모델에는 내부습도, 외부습도, 습도차, KMA 습도, 강우 flag를 함께 넣고 feature importance를 검증한다.",
            "상관이 큰 구간도 채널별로 다르므로 channel/facility 식별 feature가 필요하다.",
        ],
        Inches(7.35), Inches(1.55), Inches(5.2), Inches(4.7), 15,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "11. 모델 학습 관점의 결론", "현재 데이터만으로 최적온도=목표개폐율 기준을 확정하지 않는다")
    add_bullets(
        slide,
        [
            "목표값 정의: 관측 open_rate를 그대로 정답으로 쓰기 전에, 4~5월 상시 개방 구간을 분리하거나 가중치를 낮춘다.",
            "온도 기준: 내부온도 최소/최대 관측 범위와 2~3월 후보 범위는 제시 가능하지만, 최적온도는 운영 지식과 추가 데이터로 검증한다.",
            "Feature 설계: 내부/외부 온도, 온도차, 습도, 습도차, 강우, 일사/일조, 시간대, 계절성, lag/rolling feature를 포함한다.",
            "검증 방식: 랜덤 분할이 아니라 시간 순서 기반 walk-forward split으로 미래 누수를 방지한다.",
            "학습 전략: 회귀(open_rate 예측)와 변화 이벤트 분류(open/close 발생 여부)를 병행하는 방식을 검토한다.",
        ],
        Inches(0.85), Inches(1.45), Inches(11.85), Inches(4.85), 18,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "12. 산출물", "재현 가능한 스크립트와 발표 파일")
    add_bullets(
        slide,
        [
            "C 스크립트: data_inspect/scripts/generate_weekly_motor_weather_analysis_260211-260511.py",
            "D 스크립트: data_inspect/scripts/generate_weekly_motor_humidity_analysis_260211-260511.py",
            "PPT 생성 스크립트: data_inspect/scripts/build_data_inspect_presentation_260211-260511.py",
            f"PPT 파일: {PPTX_PATH.relative_to(PROJECT_ROOT)}",
            "요약 차트/표: data_inspect/output/presentation_260211-260511/plots/",
        ],
        Inches(0.85), Inches(1.55), Inches(11.8), Inches(4.4), 16,
    )
    add_textbox(slide, "핵심 발표 메시지: 내부온도 기준 추출은 가능하나, 4~5월 상시 개방 데이터 때문에 목표 개폐율 라벨을 그대로 학습시키면 왜곡 위험이 크다. 따라서 상세 분석 결과를 반영해 데이터 구간 분리와 feature 설계를 먼저 수행해야 한다.", Inches(0.9), Inches(6.05), Inches(11.6), Inches(0.7), 16, True, RED, PP_ALIGN.CENTER)

    prs.save(PPTX_PATH)


def main() -> None:
    summary = calculate_summary()
    plots = create_summary_plots(summary)
    make_deck(summary, plots)
    print(PPTX_PATH)


if __name__ == "__main__":
    main()
